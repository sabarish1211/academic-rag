import io
from pathlib import Path

import fitz  # PyMuPDF
from pptx import Presentation


def parse_pdf(file_bytes: bytes) -> list[dict]:
    pages = []
    doc = fitz.open(stream=file_bytes, filetype="pdf")

    for page_num, page in enumerate(doc, start=1):
        text = page.get_text("text").strip()
        if text:
            pages.append({"page": page_num, "text": text})

    doc.close()
    return pages


def parse_pptx(file_bytes: bytes) -> list[dict]:
    slides = []
    prs = Presentation(io.BytesIO(file_bytes))

    for slide_num, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        texts.append(line)

        if texts:
            slides.append({"page": slide_num, "text": "\n".join(texts)})

    return slides


def parse_document(filename: str, file_bytes: bytes) -> list[dict]:
    ext = Path(filename).suffix.lower()

    if ext == ".pdf":
        return parse_pdf(file_bytes)
    elif ext in (".pptx", ".ppt"):
        return parse_pptx(file_bytes)
    else:
        raise ValueError(f"Unsupported file type: {ext}. Only PDF and PPTX are supported.")